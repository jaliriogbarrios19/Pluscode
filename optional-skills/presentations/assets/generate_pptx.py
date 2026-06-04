#!/usr/bin/env python3
"""
Generate PowerPoint (.pptx) presentations from structured JSON input.

Usage:
    python generate_pptx.py input.json [output.pptx]

JSON Schema — top level:
{
    "meta": {
        "title": "Presentation Title",
        "subtitle": "Optional subtitle",
        "author": "Author Name",
        "date": "2026-05-29"
    },
    "theme": {
        "primary_color": "2B579A",
        "secondary_color": "5B9BD5",
        "accent_color": "ED7D31",
        "background_color": "FFFFFF",
        "text_color": "333333",
        "font_title": "Calibri",
        "font_body": "Calibri"
    },
    "slides": [ ... ]
}

Slide types and their fields:

title:
    { "type": "title", "title": "...", "subtitle": "...", "author": "...", "date": "..." }
    Note: inherits from meta if fields omitted.

section:
    { "type": "section", "title": "Section Name", "subtitle": "optional" }

bullets:
    { "type": "bullets", "title": "Slide Title", "bullets": [
        {"text": "First point", "level": 0},
        {"text": "Sub point", "level": 1},
        {"text": "Second point", "level": 0}
    ]}

content:
    { "type": "content", "title": "Slide Title", "text": "Paragraph text.\n\nAnother paragraph." }

code:
    { "type": "code", "title": "Slide Title", "code": "def hello():\\n    print('world')", "language": "python" }

two_column:
    { "type": "two_column", "title": "Slide Title",
      "left": "Left column text",
      "right": "Right column text" }

image:
    { "type": "image", "title": "Slide Title", "image_path": "/path/to/image.png", "caption": "optional" }
    image_path can be a local file path or a URL (http/https). URLs are downloaded automatically.

mermaid:
    { "type": "mermaid", "title": "Architecture Overview", "diagram": "graph TD\\n    A-->B",
      "theme": "default" }
    Renders Mermaid diagrams via mermaid.ink API and places the PNG on the slide.
    Themes: "default", "neutral", "dark", "forest".

blank:
    { "type": "blank" }

closing:
    { "type": "closing", "title": "Thank You", "subtitle": "Questions?", "contact": "email@example.com" }

All slide types support optional "notes": "Speaker notes text" and "layout": "title-only" (skips content placeholder).
"""

import json
import sys
import base64
import tempfile
import urllib.request
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DEFAULT_THEME = {
    "primary_color": "2B579A",
    "secondary_color": "5B9BD5",
    "accent_color": "ED7D31",
    "background_color": "FFFFFF",
    "text_color": "333333",
    "font_title": "Calibri",
    "font_body": "Calibri",
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def mermaid_to_png(diagram: str, theme: str = "default") -> str:
    """Encode mermaid diagram and download PNG from mermaid.ink. Returns temp file path."""
    payload = json.dumps({
        "code": diagram,
        "mermaid": json.dumps({"theme": theme})
    })
    encoded = base64.urlsafe_b64encode(payload.encode()).decode().rstrip("=")
    url = f"https://mermaid.ink/img/{encoded}"

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.close()
    urllib.request.urlretrieve(url, tmp.name)
    return tmp.name


def download_image(url: str) -> str:
    """Download an image from a URL to a temp file. Returns temp file path."""
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.close()
    urllib.request.urlretrieve(url, tmp.name)
    return tmp.name


def set_font(run, theme: dict, is_title: bool = False, size_override: int = None):
    font_name = theme.get("font_title") if is_title else theme.get("font_body")
    run.font.name = font_name or "Calibri"
    if size_override:
        run.font.size = Pt(size_override)
    elif is_title:
        run.font.size = Pt(36)
    else:
        run.font.size = Pt(18)
    run.font.color.rgb = hex_to_rgb(theme.get("text_color", "333333"))


def add_textbox(slide, left, top, width, height, text, theme, is_title=False, size=None, alignment=PP_ALIGN.LEFT, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, theme, is_title=is_title, size_override=size)
    run.font.bold = bold
    return txBox


def add_slide_number(slide, number: int, theme: dict):
    left = SLIDE_WIDTH.inches - 0.8
    top = SLIDE_HEIGHT.inches - 0.5
    add_textbox(slide, left, top, 0.6, 0.4, str(number), theme, size=10, alignment=PP_ALIGN.RIGHT)


def add_notes(slide, notes_text: str):
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def apply_background(slide, theme: dict):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme.get("background_color", "FFFFFF"))


def add_bottom_bar(slide, theme: dict):
    left = Inches(0)
    top = SLIDE_HEIGHT - Inches(0.08)
    width = SLIDE_WIDTH
    height = Inches(0.08)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "2B579A"))
    shape.line.fill.background()


def add_top_accent(slide, theme: dict):
    left = Inches(0)
    top = Inches(0)
    width = Inches(0.08)
    height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(theme.get("accent_color", "ED7D31"))
    shape.line.fill.background()


def make_title_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    apply_background(slide, theme)

    t = slide_data.get("title") or meta.get("title", "")
    st = slide_data.get("subtitle") or meta.get("subtitle", "")
    author = slide_data.get("author") or meta.get("author", "")
    date = slide_data.get("date") or meta.get("date", "")

    add_textbox(slide, 1.5, 1.5, 10.3, 1.5, t, theme, is_title=True, size=44, bold=True, alignment=PP_ALIGN.LEFT)
    if st:
        add_textbox(slide, 1.5, 3.2, 10.3, 0.8, st, theme, size=24, alignment=PP_ALIGN.LEFT)

    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    if author or date:
        footer_text = " | ".join(filter(None, [author, date]))
        add_textbox(slide, 1.5, 6.5, 10.3, 0.5, footer_text, theme, size=14, alignment=PP_ALIGN.LEFT)

    add_notes(slide, slide_data.get("notes", ""))


def make_section_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)

    primary = hex_to_rgb(theme.get("primary_color", "2B579A"))

    left = Inches(0)
    top = Inches(2.8)
    width = SLIDE_WIDTH
    height = Inches(1.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary
    shape.line.fill.background()

    t = slide_data.get("title", "Section")
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = t
    run.font.name = theme.get("font_title", "Calibri")
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    st = slide_data.get("subtitle", "")
    if st:
        add_textbox(slide, 1.5, 4.5, 10.3, 0.8, st, theme, size=18, alignment=PP_ALIGN.CENTER)

    add_notes(slide, slide_data.get("notes", ""))


def make_bullets_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "")
    add_textbox(slide, 0.8, 0.5, 11.5, 0.9, t, theme, is_title=True, size=32, bold=True)
    add_slide_number(slide, slide_num, theme)

    bullets = slide_data.get("bullets", [])
    body_top = 1.7
    body_left = 1.2
    body_width = 10.6

    txBox = slide.shapes.add_textbox(Inches(body_left), Inches(body_top), Inches(body_width), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        text = bullet.get("text", "")
        level = bullet.get("level", 0)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = level
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        font_size = 22 if level == 0 else 18
        run.font.size = Pt(font_size)
        run.font.name = theme.get("font_body", "Calibri")
        run.font.color.rgb = hex_to_rgb(theme.get("text_color", "333333"))

    add_notes(slide, slide_data.get("notes", ""))


def make_content_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "")
    add_textbox(slide, 0.8, 0.5, 11.5, 0.9, t, theme, is_title=True, size=32, bold=True)
    add_slide_number(slide, slide_num, theme)

    text = slide_data.get("text", "")
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    paragraphs = text.split("\n")
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(20)
        run.font.name = theme.get("font_body", "Calibri")
        run.font.color.rgb = hex_to_rgb(theme.get("text_color", "333333"))

    add_notes(slide, slide_data.get("notes", ""))


def make_code_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "")
    add_textbox(slide, 0.8, 0.5, 11.5, 0.9, t, theme, is_title=True, size=32, bold=True)
    add_slide_number(slide, slide_num, theme)

    code = slide_data.get("code", "")
    language = slide_data.get("language", "")

    code_bg = hex_to_rgb("1E1E1E")
    left = Inches(0.8)
    top = Inches(1.7)
    width = Inches(11.5)
    height = Inches(5.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = code_bg
    shape.line.fill.background()

    if language:
        add_textbox(slide, 0.95, 1.72, 3, 0.35, language.upper(), theme, size=10)
        # Override color for the lang label on dark bg
        lbl = slide.shapes.add_textbox(Inches(0.95), Inches(1.72), Inches(3), Inches(0.35))
        lp = lbl.text_frame.paragraphs[0]
        lr = lp.add_run()
        lr.text = language.upper()
        lr.font.size = Pt(10)
        lr.font.name = "Consolas"
        lr.font.color.rgb = hex_to_rgb("888888")

    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.15), Inches(10.9), Inches(4.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = code.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(14)
        run.font.name = "Consolas"
        run.font.color.rgb = hex_to_rgb("D4D4D4")

    add_notes(slide, slide_data.get("notes", ""))


def make_two_column_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "")
    add_textbox(slide, 0.8, 0.5, 11.5, 0.9, t, theme, is_title=True, size=32, bold=True)
    add_slide_number(slide, slide_num, theme)

    col_width = 5.0
    col_top = 1.7
    col_height = 5.0
    gap = 0.3

    left_text = slide_data.get("left", "")
    right_text = slide_data.get("right", "")

    l_box = slide.shapes.add_textbox(Inches(1.0), Inches(col_top), Inches(col_width), Inches(col_height))
    l_box.text_frame.word_wrap = True
    lp = l_box.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_text
    lr.font.size = Pt(18)
    lr.font.name = theme.get("font_body", "Calibri")
    lr.font.color.rgb = hex_to_rgb(theme.get("text_color", "333333"))

    r_box = slide.shapes.add_textbox(Inches(1.0 + col_width + gap), Inches(col_top), Inches(col_width), Inches(col_height))
    r_box.text_frame.word_wrap = True
    rp = r_box.text_frame.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_text
    rr.font.size = Pt(18)
    rr.font.name = theme.get("font_body", "Calibri")
    rr.font.color.rgb = hex_to_rgb(theme.get("text_color", "333333"))

    divider_left = Inches(1.0 + col_width + gap / 2 - 0.005)
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, divider_left, Inches(col_top), Inches(0.01), Inches(col_height))
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(theme.get("secondary_color", "5B9BD5"))
    divider.line.fill.background()

    add_notes(slide, slide_data.get("notes", ""))


def make_image_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "")
    add_textbox(slide, 0.8, 0.5, 11.5, 0.9, t, theme, is_title=True, size=32, bold=True)
    add_slide_number(slide, slide_num, theme)

    image_path = slide_data.get("image_path", "")
    caption = slide_data.get("caption", "")

    picture_path = None
    if image_path:
        if image_path.startswith(("http://", "https://")):
            try:
                picture_path = download_image(image_path)
            except Exception as e:
                print(f"Warning: failed to download image from {image_path}: {e}")
        elif Path(image_path).exists():
            picture_path = image_path

    if picture_path:
        slide.shapes.add_picture(picture_path, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.8))
    else:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.8))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb("E0E0E0")
        placeholder.line.color.rgb = hex_to_rgb("CCCCCC")
        tf = placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "[Image — provide a local path or URL]" + (f"\n{caption}" if caption else "")
        run.font.size = Pt(16)
        run.font.color.rgb = hex_to_rgb("888888")

    if caption and picture_path:
        add_textbox(slide, 1.5, 6.7, 10.3, 0.4, caption, theme, size=12, alignment=PP_ALIGN.CENTER)

    add_notes(slide, slide_data.get("notes", ""))


def make_mermaid_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "")
    add_textbox(slide, 0.8, 0.5, 11.5, 0.9, t, theme, is_title=True, size=32, bold=True)
    add_slide_number(slide, slide_num, theme)

    diagram = slide_data.get("diagram", "")
    mermaid_theme = slide_data.get("theme", "default")
    caption = slide_data.get("caption", "")

    if diagram:
        try:
            png_path = mermaid_to_png(diagram, mermaid_theme)
            slide.shapes.add_picture(png_path, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.2))
        except Exception as e:
            err_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.2))
            err_box.fill.solid()
            err_box.fill.fore_color.rgb = hex_to_rgb("FFF3F3")
            err_box.line.color.rgb = hex_to_rgb("FF4444")
            tf = err_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"Mermaid render failed: {e}\n\nDiagram:\n{diagram[:200]}"
            run.font.size = Pt(12)
            run.font.color.rgb = hex_to_rgb("CC0000")

    if caption:
        add_textbox(slide, 0.8, 7.0, 11.5, 0.3, caption, theme, size=11, alignment=PP_ALIGN.CENTER)

    add_notes(slide, slide_data.get("notes", ""))


def make_blank_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)
    add_slide_number(slide, slide_num, theme)
    add_notes(slide, slide_data.get("notes", ""))


def make_closing_slide(prs, slide_data: dict, meta: dict, theme: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, theme)
    add_top_accent(slide, theme)
    add_bottom_bar(slide, theme)

    t = slide_data.get("title", "Thank You")
    add_textbox(slide, 1.5, 2.0, 10.3, 1.2, t, theme, is_title=True, size=48, bold=True, alignment=PP_ALIGN.CENTER)

    st = slide_data.get("subtitle", "")
    if st:
        add_textbox(slide, 1.5, 3.5, 10.3, 0.8, st, theme, size=24, alignment=PP_ALIGN.CENTER)

    contact = slide_data.get("contact", "")
    if contact:
        add_textbox(slide, 1.5, 5.0, 10.3, 0.6, contact, theme, size=18, alignment=PP_ALIGN.CENTER)

    add_notes(slide, slide_data.get("notes", ""))


SLIDE_BUILDERS = {
    "title": make_title_slide,
    "section": make_section_slide,
    "bullets": make_bullets_slide,
    "content": make_content_slide,
    "code": make_code_slide,
    "two_column": make_two_column_slide,
    "image": make_image_slide,
    "mermaid": make_mermaid_slide,
    "blank": make_blank_slide,
    "closing": make_closing_slide,
}


def generate_pptx(input_data: dict, output_path: str) -> str:
    meta = input_data.get("meta", {})
    theme = {**DEFAULT_THEME, **input_data.get("theme", {})}
    slides = input_data.get("slides", [])

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, slide_data in enumerate(slides):
        slide_type = slide_data.get("type", "bullets")
        builder = SLIDE_BUILDERS.get(slide_type)
        if builder:
            builder(prs, slide_data, meta, theme, i + 1)
        else:
            print(f"Warning: unknown slide type '{slide_type}' at index {i}, skipping")

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("presentation.pptx")

    if not input_path.exists():
        print(f"Error: input file not found: {input_path}")
        sys.exit(1)

    with open(input_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    result = generate_pptx(data, str(output_path))
    print(f"Presentation saved to: {result}")
    print(f"Slides: {len(data.get('slides', []))}")


if __name__ == "__main__":
    main()
